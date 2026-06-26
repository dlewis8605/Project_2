import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme colors
    bg_color = RGBColor(11, 8, 22)         # #0b0816 (Deep Purple Black)
    primary_color = RGBColor(168, 85, 247) # #a855f7 (Purple Accent)
    accent_color = RGBColor(236, 72, 153)  # #ec4899 (Pink Accent)
    text_light = RGBColor(243, 241, 254)   # #f3f1fe (Off White)
    text_muted = RGBColor(159, 155, 191)   # #9f9bbf (Light Gray Purple)
    card_bg = RGBColor(20, 16, 38)         # #141026 (Translucent Glass Card)
    
    # Helper: Set slide dark background
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Helper: Add Standard Title
    def add_slide_header(slide, title_text, category_text=None):
        set_background(slide)
        
        # Category label (tiny header)
        if category_text:
            catBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.4))
            cat_tf = catBox.text_frame
            cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
            p_cat = cat_tf.paragraphs[0]
            p_cat.text = category_text.upper()
            p_cat.font.name = 'Calibri'
            p_cat.font.size = Pt(12)
            p_cat.font.bold = True
            p_cat.font.color.rgb = accent_color
            
        # Main title
        titleBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.8), Inches(1.0))
        tf = titleBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = primary_color

    # Slide 1: Cover Slide
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1)
    
    # Decorative visual: Glowing accent band
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary_color
    shape.line.fill.background()
    
    # Text container
    coverBox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.5), Inches(4.5))
    tf = coverBox.text_frame
    tf.word_wrap = True
    
    # Project Title
    p_title = tf.paragraphs[0]
    p_title.text = "AURA — UX/UI Asset Catalog"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(48)
    p_title.font.bold = True
    p_title.font.color.rgb = primary_color
    p_title.space_after = Pt(10)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "A Curated Digital Asset Repository for Web Creators"
    p_sub.font.name = 'Calibri'
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = text_light
    p_sub.space_after = Pt(40)
    
    # Metadata
    p_meta = tf.add_paragraph()
    p_meta.text = "Student Name: Darell Lewis\nCourse: Web Development II (Back-End)\nInstructor: Westcliff Course Reviewer\nDate: June 2026"
    p_meta.font.name = 'Calibri'
    p_meta.font.size = Pt(16)
    p_meta.font.color.rgb = text_muted
    p_meta.line_spacing = 1.3

    # Slide 2: Project Introduction
    slide2 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide2, "Introduction & Project Goals", "Overview")
    
    bodyBox = slide2.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf2 = bodyBox.text_frame
    tf2.word_wrap = True
    
    bullets = [
        ("The Challenge:", "Web developers waste countless hours reinventing common UI design elements (buttons, glassmorphic panels, gradients, color palettes)."),
        ("The Solution (AURA):", "A secure, lightning-fast database-driven asset catalog that lets creators browse, preview, copy, and save curated web assets."),
        ("Key Goals & Objectives:", ""),
        ("  • Centralized Catalog:", "Efficiently store and organize HTML/CSS/SASS and SVG assets by categories and tags."),
        ("  • User Personalization:", "Allow developers to register, sign in, and maintain a curated dashboard of favorited design snippets."),
        ("  • Robust Data Integrity:", "Facilitate community contributions through strict server-side validation and sanitization."),
        ("  • Performance & Portability:", "Utilize Progressive Web App (PWA) strategies for offline catalog access and seamless installation.")
    ]
    
    first = True
    for header, desc in bullets:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(8)
        
        if header.strip().startswith("•") or header.strip() == "":
            p.font.size = Pt(18)
            run1 = p.add_run()
            run1.text = header
            run1.font.bold = True
            run1.font.color.rgb = accent_color
        else:
            p.font.size = Pt(20)
            run1 = p.add_run()
            run1.text = header + " "
            run1.font.bold = True
            run1.font.color.rgb = text_light
            
        if desc:
            run2 = p.add_run()
            run2.text = desc
            run2.font.color.rgb = text_muted

    # Slide 3: Technology Stack
    slide3 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide3, "Modern Stack & Architecture", "Tech Stack")
    
    # 3 columns for Front-End, Back-End, Database/DevOps
    col_width = Inches(3.6)
    col_gap = Inches(0.5)
    left_start = Inches(0.75)
    top_pos = Inches(2.2)
    height = Inches(4.3)
    
    techs = [
        ("FRONT-END", [
            "• Vue.js CDN (Client Engine)",
            "  Handles reactive search filter & copy events instantly",
            "• SASS Custom Variables",
            "  Provides glassmorphism tokens and radial orb animations",
            "• Bootstrap 4 Framework",
            "  Ensures fluid responsive grids across devices"
        ], primary_color),
        ("BACK-END RUNTIME", [
            "• Node.js & Express.js",
            "  Lightweight, asynchronous event-driven server runtime",
            "• express-session & cookies",
            "  Maintains stateful user credentials securely",
            "• BCryptJS Hashing",
            "  Protects stored user passwords with cryptographic salting"
        ], accent_color),
        ("DATABASE & SECURITY", [
            "• MongoDB (NoSQL Database)",
            "  Document store tailored for dynamic schema objects",
            "• Mongoose ODM",
            "  Declares schema constraints & data integrity models",
            "• OWASP Middleware",
            "  Helmet security headers, rate limiting, and mongo-sanitizer"
        ], text_light)
    ]
    
    for idx, (title, points, color) in enumerate(techs):
        col_left = left_start + idx * (col_width + col_gap)
        
        # Draw background card shape
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_left, top_pos, col_width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = primary_color if idx == 0 else (accent_color if idx == 1 else text_muted)
        card.line.width = Pt(1.5)
        
        # Add textbox inside
        box = slide3.shapes.add_textbox(col_left + Inches(0.2), top_pos + Inches(0.2), col_width - Inches(0.4), height - Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_header = tf.paragraphs[0]
        p_header.text = title
        p_header.font.name = 'Arial'
        p_header.font.size = Pt(20)
        p_header.font.bold = True
        p_header.font.color.rgb = color
        p_header.space_after = Pt(15)
        
        for pt in points:
            p = tf.add_paragraph()
            p.text = pt
            p.font.name = 'Calibri'
            if pt.strip().startswith("•"):
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = text_light
                p.space_before = Pt(8)
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = text_muted
                p.space_after = Pt(4)

    # Slide 4: Database Schemas
    slide4 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide4, "Mongoose Database Models", "Data Structure")
    
    # 2 columns: User Model & Asset Model
    col_width = Inches(5.6)
    col_gap = Inches(0.6)
    left_start = Inches(0.75)
    top_pos = Inches(2.0)
    height = Inches(4.6)
    
    models = [
        ("User Model Schema", [
            "Stores details of creators & catalog members.",
            "• username: String (Unique, required, min 5 chars)",
            "• email: String (Unique, required, lowercased, regex validated)",
            "• password: String (Hashed via bcrypt salt values)",
            "• favorites: [ ObjectId ] (Array referencing Asset documents)",
            "• createdAt: Date (Auto-generated timestamp)"
        ]),
        ("Asset Model Schema", [
            "Houses metadata & code contents for snippets.",
            "• title: String (Required, trimmed, min 3 chars)",
            "• description: String (Required text explanation)",
            "• category: String (Enum: css, sass, palette, svg, other)",
            "• code: String (Raw code snippet text)",
            "• tags: [ String ] (Array of search tokens)",
            "• creator: ObjectId (Refers back to contributing User)",
            "• downloads: Number (Stat tracker counter for copy events)"
        ])
    ]
    
    for idx, (title, points) in enumerate(models):
        col_left = left_start + idx * (col_width + col_gap)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_left, top_pos, col_width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = primary_color if idx == 0 else accent_color
        card.line.width = Pt(1.5)
        
        box = slide4.shapes.add_textbox(col_left + Inches(0.2), top_pos + Inches(0.2), col_width - Inches(0.4), height - Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_header = tf.paragraphs[0]
        p_header.text = title
        p_header.font.name = 'Arial'
        p_header.font.size = Pt(22)
        p_header.font.bold = True
        p_header.font.color.rgb = primary_color if idx == 0 else accent_color
        p_header.space_after = Pt(12)
        
        for pt in points:
            p = tf.add_paragraph()
            p.text = pt
            p.font.name = 'Calibri'
            if pt.strip().startswith("•"):
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = text_light
                p.space_before = Pt(6)
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = text_muted
                p.space_after = Pt(3)

    # Slide 5: User Auth & Session Management
    slide5 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide5, "Secure Authentication & Sessions", "Security")
    
    bodyBox = slide5.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf5 = bodyBox.text_frame
    tf5.word_wrap = True
    
    auth_points = [
        ("Password Hashing (BCrypt):", "User passwords are encrypted before database insertion using genSalt and hash. The raw password is never stored or compared in plain text."),
        ("Session Management:", "Implemented stateful express-session middleware. Sessions are securely signed with a cryptographic secret key."),
        ("Persistent Session Store (MongoStore):", "Instead of using leaky in-memory stores, session states are written to the MongoDB sessions collection, ensuring logins persist across server restarts."),
        ("Security Cookies:", "Cookies are configured with httpOnly: true to prevent cross-site scripting (XSS) client-side access. Session timeouts are set to expire automatically in 24 hours."),
        ("Session Check & Protected API Routes:", "A validation middleware requireAuth intercepts asset submission and favoriting requests, returning a 401 Unauthorized status if the active session ID is missing or invalid.")
    ]
    
    first = True
    for header, desc in auth_points:
        p = tf5.paragraphs[0] if first else tf5.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.font.size = Pt(18)
        
        run1 = p.add_run()
        run1.text = header + " "
        run1.font.bold = True
        run1.font.color.rgb = accent_color
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.color.rgb = text_muted

    # Slide 6: PWA Integration
    slide6 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide6, "PWA & Offline-First Implementation", "Offline Capability")
    
    bodyBox = slide6.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf6 = bodyBox.text_frame
    tf6.word_wrap = True
    
    pwa_points = [
        ("Web App Manifest (manifest.json):", "Defines the app's standalone desktop/mobile UI metadata, background color, start URL, and lists high-res maskable icon files (192px and 512px)."),
        ("Service Worker Registration (pwa.js):", "Registers the background worker thread on browser startup, handles system offline state triggers, and renders a glowing banner to inform the user of offline status."),
        ("Cache Lifecycle Installation (sw.js):", "Pre-caches static application assets (HTML, CSS, JS, manifest, icons, and Bootstrap/FontAwesome CDN libraries) upon install, creating an instantly-loading app shell."),
        ("Caching Interception Strategy:", "Implements a custom fetch caching listener. Static assets are retrieved utilizing a Cache-First approach for speed, while API endpoints default to Network-Only with an offline JSON fallback."),
        ("Offline Experience:", "Users can disconnect from the internet and still launch AURA, browse cached code blocks, and read previously fetched CSS components seamlessly.")
    ]
    
    first = True
    for header, desc in pwa_points:
        p = tf6.paragraphs[0] if first else tf6.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.font.size = Pt(18)
        
        run1 = p.add_run()
        run1.text = header + " "
        run1.font.bold = True
        run1.font.color.rgb = primary_color
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.color.rgb = text_muted

    # Slide 7: OWASP Security Middlewares
    slide7 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide7, "OWASP Top-10 Defenses & Protections", "Security Hardening")
    
    bodyBox = slide7.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf7 = bodyBox.text_frame
    tf7.word_wrap = True
    
    security_points = [
        ("HTTP Security Headers (Helmet):", "Express utilizes Helmet to set secure response headers, disabling power leaks (X-Powered-By) and restricting Content Security Policies (CSP) to allow only verified CDNs."),
        ("NoSQL Query Injection Prevention (express-mongo-sanitize):", "Sanitizes user input properties by recursively stripping symbols like '$' and '.' to prevent query modification attacks on MongoDB."),
        ("Cross-Site Scripting (XSS) Defenses (xss-clean):", "Sanitizes request body parameters, queries, and parameters to filter out malicious HTML/script injections, neutralizing DOM manipulation attempts."),
        ("API Rate Limiting & Brute-Force Blocker (express-rate-limit):", "Sets maximum requests per IP window. Protects authentication paths (/api/auth/login) with a tighter rate limit (30 requests/15m) to stop automated brute-force attacks."),
        ("Request Size Limiting:", "Body parsers restrict JSON and URL payload parameters to 10kb maximum to mitigate Denial of Service (DoS) memory exhaustion vectors.")
    ]
    
    first = True
    for header, desc in security_points:
        p = tf7.paragraphs[0] if first else tf7.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.font.size = Pt(18)
        
        run1 = p.add_run()
        run1.text = header + " "
        run1.font.bold = True
        run1.font.color.rgb = accent_color
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.color.rgb = text_muted

    # Slide 8: Automated Verification & Unit Tests
    slide8 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide8, "Quality Assurance: Unit & API Testing", "Verification")
    
    bodyBox = slide8.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf8 = bodyBox.text_frame
    tf8.word_wrap = True
    
    test_points = [
        ("Testing Suite (Jest & Supertest):", "Leverages Jest test runner along with Supertest to spin up ephemeral Express server routes and query endpoints without bindings to active ports."),
        ("Schema Validation Assertions:", "Tests User registration constraints (username minimum length, email validation regex, duplicate checking) and verifies mongoose schema violations return correct error lists."),
        ("API Endpoint Tests:", "Validates authorization states: GET /api/assets should load for anyone, but POST /api/assets (asset submission) must fail with a 415/401 status code if the user is unauthenticated."),
        ("Database Isolation:", "Configures isolated test hooks (beforeAll, afterAll) to connect to the database, seed test documents, clean state variables, and disconnect gracefully after validation completes."),
        ("Continuous Integration Ready:", "Defines clean NPM test scripts running Jest in single-band detection mode for immediate pipeline alignment.")
    ]
    
    first = True
    for header, desc in test_points:
        p = tf8.paragraphs[0] if first else tf8.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.font.size = Pt(18)
        
        run1 = p.add_run()
        run1.text = header + " "
        run1.font.bold = True
        run1.font.color.rgb = primary_color
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.color.rgb = text_muted

    # Slide 9: Demo Steps & Walkthrough
    slide9 = prs.slides.add_slide(slide_layout)
    add_slide_header(slide9, "Interactive System Walkthrough", "Live Demo")
    
    bodyBox = slide9.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
    tf9 = bodyBox.text_frame
    tf9.word_wrap = True
    
    demo_steps = [
        ("Step 1: Start Database & Server", "Ensure MongoDB is running, then run `npm start` to spin up the local server on http://localhost:3000."),
        ("Step 2: Guest Experience", "Browse the catalog dashboard, search for keyword '#glassmorphism', toggle categories (CSS, SASS), copy source codes, and observe the download count incrementing."),
        ("Step 3: User Authentication", "Navigate to registration, enter credentials, and create an account. Watch the UI immediately update the navbar options and greet the user."),
        ("Step 4: Contribute to Catalog", "Go to the Asset Submission form, enter details for a custom neon glass-button styling, input tags, and hit Publish. View it load immediately at the top of the home page."),
        ("Step 5: Curate Dashboard & Offline Mode", "Favorite the new asset, visit the Favorites dashboard tab, then simulate network disconnection. Observe the offline warning and confirm that you can still browse previously fetched assets.")
    ]
    
    first = True
    for header, desc in demo_steps:
        p = tf9.paragraphs[0] if first else tf9.add_paragraph()
        first = False
        p.space_after = Pt(10)
        p.font.size = Pt(18)
        
        run1 = p.add_run()
        run1.text = header + " — "
        run1.font.bold = True
        run1.font.color.rgb = accent_color
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.color.rgb = text_muted

    # Save presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"PowerPoint Presentation successfully generated as {output_filename}")

if __name__ == '__main__':
    create_presentation()
